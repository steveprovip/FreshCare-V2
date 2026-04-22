from docx import Document
from docx.shared import Cm, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
import glob
import os

def update_word():
    print("Updating Word with python-docx...")
    doc = Document("BaoCaoCongNghePhanMem_Nhom8_KTKN1_N08_V11_Final.docx")

    # Set page margins
    for section in doc.sections:
        section.page_height = Cm(29.7)
        section.page_width = Cm(21.0)
        section.top_margin = Cm(2.0)
        section.bottom_margin = Cm(2.0)
        section.left_margin = Cm(3.0)
        section.right_margin = Cm(1.5)

    # Set Normal style
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = Pt(14)
    style.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    style.paragraph_format.line_spacing = 1.5

    def replace_text_in_runs(paragraph):
        # We need to assemble the full text, replace, and clear old runs then add a new run
        # because replacing across runs in python-docx is messy. We just replace text if found.
        # But for simple words, replacing in paragraph.text strips inline formatting like bold.
        # However, for this document, inline formatting on the specific strings "3 ngày" might not be crucial,
        # or we just iterate runs if the text is in a single run.
        text = paragraph.text
        if "<= 3 ngày" in text or "< 3 ngày" in text or "<= 3" in text:
            # We will use simple replace if text is simple
            for run in paragraph.runs:
                if "<= 3 ngày" in run.text:
                    run.text = run.text.replace("<= 3 ngày", "<= 14 ngày")
                if "< 3 ngày" in run.text:
                    run.text = run.text.replace("< 3 ngày", "<= 14 ngày")
                if "<= 3" in run.text:
                    run.text = run.text.replace("<= 3", "<= 14")
            
            # If the word spans multiple runs, fallback to paragraph.text assignment
            if ("<= 3" in paragraph.text) or ("< 3 ngày" in paragraph.text):
                # Fallback
                new_text = paragraph.text.replace("<= 3 ngày", "<= 14 ngày").replace("< 3 ngày", "<= 14 ngày").replace("<= 3", "<= 14")
                paragraph.text = new_text

    for paragraph in doc.paragraphs:
        replace_text_in_runs(paragraph)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    replace_text_in_runs(paragraph)

    doc.save("BaoCaoCongNghePhanMem_Nhom8_KTKN1_N08_V11_Final.docx")
    print("Word updated.")

def update_ppt():
    print("Updating PPT with python-pptx...")
    ppts = glob.glob("..\Slide*.pptx")
    if not ppts:
        print("No PPT found")
        return
    ppt_path = ppts[0]
    prs = Presentation(ppt_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "<= 3 ngày" in run.text:
                            run.text = run.text.replace("<= 3 ngày", "<= 14 ngày")
                        if "< 3 ngày" in run.text:
                            run.text = run.text.replace("< 3 ngày", "<= 14 ngày")
                        if "<= 3" in run.text:
                            run.text = run.text.replace("<= 3", "<= 14")
                    
                    # check if still present due to cross-run
                    if "<= 3" in paragraph.text:
                         # replace all text
                         t = paragraph.text.replace("<= 3 ngày", "<= 14 ngày").replace("< 3 ngày", "<= 14 ngày").replace("<= 3", "<= 14")
                         paragraph.clear() # removes all runs
                         p = paragraph.add_run()
                         p.text = t

    prs.save(ppt_path)
    print("PPT updated.")

if __name__ == "__main__":
    update_word()
    update_ppt()
